"""Module for extracting text from files of various formats."""

import asyncio
import concurrent.futures
import io
import logging
import os
import shutil
import subprocess
import tarfile
import tempfile
import threading
import time
import xml.etree.ElementTree as ET
import zipfile
from pathlib import Path
from typing import Any, Dict, List, Optional

# Archive imports
try:
    import rarfile
except ImportError:
    rarfile = None

try:
    import py7zr
except ImportError:
    py7zr = None

# Imports for various formats
try:
    import pdfplumber
    import PyPDF2
except ImportError:
    PyPDF2 = None
    pdfplumber = None

try:
    from docx import Document
except ImportError:
    Document = None

try:
    import pandas as pd
except ImportError:
    pd = None

try:
    import pytesseract
    from PIL import Image
except ImportError:
    Image = None
    pytesseract = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    from bs4 import BeautifulSoup
except ImportError:
    BeautifulSoup = None

try:
    import markdown
except ImportError:
    markdown = None

try:
    from odf.opendocument import load
    from odf.teletype import extractText
    from odf.text import P
except ImportError:
    load = None
    P = None
    extractText = None

try:
    from striprtf.striprtf import rtf_to_text
except ImportError:
    rtf_to_text = None

try:
    import yaml
except ImportError:
    yaml = None

# Web extraction (new in v1.10.0)
try:
    import ipaddress
    from urllib.parse import urljoin, urlparse

    import requests
except ImportError:
    requests = None
    urljoin = None
    urlparse = None
    ipaddress = None

# Playwright for JS rendering (new in v1.10.1)
try:
    from playwright.sync_api import sync_playwright
except ImportError:
    sync_playwright = None

from fastapi import BackgroundTasks

from app.config import settings
from app.utils import get_file_extension, is_archive_format, is_supported_format

logger = logging.getLogger(__name__)


class TextExtractor:
    """Class for extracting text from files of various formats."""

    def __init__(self):
        """Initialize the text extractor."""
        self.ocr_languages = settings.OCR_LANGUAGES
        self.timeout = settings.PROCESSING_TIMEOUT_SECONDS
        # Create thread pool for CPU-bound operations
        self._thread_pool = concurrent.futures.ThreadPoolExecutor(max_workers=4)

    def extract_text(self, file_content: bytes, filename: str) -> List[Dict[str, Any]]:
        """Main text extraction method (synchronous for execution in threadpool)."""
        # Check if file is an archive
        if is_archive_format(filename, settings.SUPPORTED_FORMATS):
            return self._extract_from_archive(file_content, filename)

        # Check format support
        if not is_supported_format(filename, settings.SUPPORTED_FORMATS):
            raise ValueError(f"Unsupported file format: {filename}")

        # Check MIME type for security (synchronous operation)
        is_valid_mime = self._check_mime_type(file_content, filename)

        if not is_valid_mime:
            logger.warning(f"File {filename} MIME type does not match its extension")
            # Don't block, but warn

        extension = get_file_extension(filename)

        # Check extension is not None
        if not extension:
            raise ValueError(f"Could not determine file extension for: {filename}")

        try:
            # Extract text synchronously
            text = self._extract_text_by_format(file_content, extension, filename)

            # Return array with one element for consistency
            return [
                {
                    "filename": filename,
                    "path": filename,
                    "size": len(file_content),
                    "type": extension,
                    "text": text.strip() if text else "",
                }
            ]

        except Exception as e:
            logger.error(f"Error extracting text from {filename}: {str(e)}")
            raise ValueError(f"Error extracting text: {str(e)}")

    def _extract_text_by_format(
        self, content: bytes, extension: str, filename: str
    ) -> str:
        """Extract text depending on format (synchronous version)."""
        # Create dictionary mapping extensions to extraction methods
        extraction_methods = self._get_extraction_methods_mapping()

        # Check if file is source code
        source_code_extensions = settings.SUPPORTED_FORMATS.get("source_code", [])
        if extension in source_code_extensions:
            return self._extract_from_source_code_sync(content, extension, filename)

        # Look for suitable extraction method
        extractor_method = extraction_methods.get(extension)
        if extractor_method:
            return extractor_method(content)

        # Check extension groups
        for extensions_group, method in self._get_group_extraction_methods():
            if extension in extensions_group:
                return method(content)

        raise ValueError(f"Unsupported file format: {extension}")

    def _get_extraction_methods_mapping(self) -> dict:
        """Get dictionary mapping extensions to extraction methods."""
        return {
            "pdf": self._extract_from_pdf_sync,
            "docx": self._extract_from_docx_sync,
            "doc": self._extract_from_doc_sync,
            "csv": self._extract_from_csv_sync,
            "pptx": self._extract_from_pptx_sync,
            "ppt": self._extract_from_ppt_sync,
            "txt": self._extract_from_txt_sync,
            "json": self._extract_from_json_sync,
            "rtf": self._extract_from_rtf_sync,
            "odt": self._extract_from_odt_sync,
            "xml": self._extract_from_xml_sync,
            "epub": self._extract_from_epub_sync,
            "eml": self._extract_from_eml_sync,
            "msg": self._extract_from_msg_sync,
        }

    def _get_group_extraction_methods(self) -> list:
        """Get list of extension groups with corresponding methods."""
        return [
            (["xls", "xlsx"], self._extract_from_excel_sync),
            (
                ["jpg", "jpeg", "png", "tiff", "tif", "bmp", "gif"],
                self._extract_from_image_sync,
            ),
            (["html", "htm"], self._extract_from_html_sync),
            (["md", "markdown"], self._extract_from_markdown_sync),
            (["yaml", "yml"], self._extract_from_yaml_sync),
        ]

    def _extract_from_pdf_sync(self, content: bytes) -> str:
        """Synchronous text extraction from PDF."""
        if not pdfplumber:
            raise ImportError("pdfplumber is not installed")

        text_parts = []
        temp_file_path = None

        try:
            with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as temp_file:
                temp_file.write(content)
                temp_file_path = temp_file.name

            with pdfplumber.open(temp_file_path) as pdf:
                for page_num, page in enumerate(pdf.pages, 1):
                    page_texts = self._extract_pdf_page_content(page, page_num)
                    text_parts.extend(page_texts)

            return "\n\n".join(text_parts)

        except Exception as e:
            logger.error(f"Error processing PDF: {str(e)}")
            raise ValueError(f"Error processing PDF: {str(e)}")
        finally:
            self._cleanup_temp_file(temp_file_path)

    def _extract_pdf_page_content(self, page, page_num: int) -> list:
        """Extract PDF page content."""
        page_texts = []

        # Extract text from page
        page_text = page.extract_text()
        if page_text:
            page_texts.append(f"[Page {page_num}]\n{page_text}")

        # Extract images and OCR
        if page.images:
            page_texts.extend(self._extract_pdf_page_images(page))

        return page_texts

    def _extract_pdf_page_images(self, page) -> list:
        """Extract text from images on PDF page."""
        image_texts = []

        for img_idx, img in enumerate(page.images):
            try:
                image_text = self._ocr_from_pdf_image_sync(page, img)
                if image_text.strip():
                    image_texts.append(f"[Image {img_idx + 1}]\n{image_text}")
            except Exception as e:
                logger.warning(f"OCR error on image {img_idx + 1}: {str(e)}")

        return image_texts

    def _cleanup_temp_file(self, temp_file_path: str) -> None:
        """Safely remove temporary file."""
        if temp_file_path and os.path.exists(temp_file_path):
            try:
                os.unlink(temp_file_path)
            except OSError as e:
                logger.warning(
                    f"Failed to remove temporary file {temp_file_path}: {str(e)}"
                )

    def _extract_from_docx_sync(self, content: bytes) -> str:
        """Synchronous text extraction from DOCX with full extraction according to project specs."""
        if not Document:
            raise ImportError("python-docx is not installed")

        try:
            doc = Document(io.BytesIO(content))
            text_parts = []

            # Extract various parts of the document
            text_parts.extend(self._extract_docx_paragraphs(doc))
            text_parts.extend(self._extract_docx_tables(doc))
            text_parts.extend(self._extract_docx_headers_footers(doc))
            text_parts.extend(self._extract_docx_footnotes(doc))
            text_parts.extend(self._extract_docx_comments(doc))

            return "\n\n".join(text_parts)

        except Exception as e:
            logger.error(f"Error processing DOCX: {str(e)}")
            raise ValueError(f"Error processing DOCX: {str(e)}")

    def _extract_docx_paragraphs(self, doc) -> list:
        """Extract main text from DOCX paragraphs."""
        text_parts = []
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_parts.append(paragraph.text)
        return text_parts

    def _extract_docx_tables(self, doc) -> list:
        """Extract text from DOCX tables."""
        text_parts = []
        for table in doc.tables:
            table_text = []
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    row_text.append(cell.text.strip())
                table_text.append("\t".join(row_text))

            if table_text:
                text_parts.append("\n".join(table_text))
        return text_parts

    def _extract_docx_headers_footers(self, doc) -> list:
        """Extract text from DOCX headers and footers."""
        text_parts = []
        for section in doc.sections:
            # Extract header
            if section.header:
                header_text = self._extract_section_text(section.header.paragraphs)
                if header_text:
                    text_parts.append(
                        f"[Header]\n{' '.join(header_text)}"
                    )

            # Extract footer
            if section.footer:
                footer_text = self._extract_section_text(section.footer.paragraphs)
                if footer_text:
                    text_parts.append(f"[Footer]\n{' '.join(footer_text)}")
        return text_parts

    def _extract_section_text(self, paragraphs) -> list:
        """Extract text from section paragraphs."""
        text_parts = []
        for paragraph in paragraphs:
            if paragraph.text.strip():
                text_parts.append(paragraph.text)
        return text_parts

    def _extract_docx_footnotes(self, doc) -> list:
        """Extract footnotes from DOCX."""
        text_parts = []
        try:
            if hasattr(doc, "footnotes") and doc.footnotes:
                footnotes_text = []
                for footnote in doc.footnotes:
                    if hasattr(footnote, "paragraphs"):
                        footnote_text = self._extract_section_text(footnote.paragraphs)
                        footnotes_text.extend(footnote_text)
                if footnotes_text:
                    text_parts.append(f"[Footnotes]\n{' '.join(footnotes_text)}")
        except Exception as e:
            logger.debug(f"Failed to extract footnotes from DOCX: {str(e)}")
        return text_parts

    def _extract_docx_comments(self, doc) -> list:
        """Extract comments from DOCX."""
        text_parts = []
        try:
            if hasattr(doc, "comments") and doc.comments:
                comments_text = []
                for comment in doc.comments:
                    if hasattr(comment, "paragraphs"):
                        comment_text = self._extract_section_text(comment.paragraphs)
                        comments_text.extend(comment_text)
                if comments_text:
                    text_parts.append(f"[Comments]\n{' '.join(comments_text)}")
        except Exception as e:
            logger.debug(f"Failed to extract comments from DOCX: {str(e)}")
        return text_parts

    def _extract_from_doc_sync(self, content: bytes) -> str:
        """Synchronous text extraction from DOC via conversion to DOCX using LibreOffice."""
        if not Document:
            raise ImportError("python-docx is not installed")

        try:
            # Create temporary files
            with tempfile.NamedTemporaryFile(suffix=".doc", delete=False) as temp_doc:
                temp_doc.write(content)
                temp_doc_path = temp_doc.name

            # Create temporary output directory
            temp_dir = tempfile.mkdtemp()

            try:
                # Convert .doc to .docx using LibreOffice with resource limits
                from .config import settings
                from .utils import run_subprocess_with_limits

                result = run_subprocess_with_limits(
                    command=[
                        "libreoffice",
                        "--headless",
                        "--convert-to",
                        "docx",
                        "--outdir",
                        temp_dir,
                        temp_doc_path,
                    ],
                    timeout=30,
                    memory_limit=settings.MAX_LIBREOFFICE_MEMORY,
                    capture_output=True,
                    text=True,
                )

                if result.returncode != 0:
                    logger.error(f"LibreOffice conversion failed: {result.stderr}")
                    raise ValueError("Failed to convert DOC to DOCX")

                # Find converted file
                doc_filename = os.path.splitext(os.path.basename(temp_doc_path))[0]
                docx_path = os.path.join(temp_dir, f"{doc_filename}.docx")

                if not os.path.exists(docx_path):
                    raise ValueError("Converted DOCX file not found")

                # Read converted DOCX file
                with open(docx_path, "rb") as docx_file:
                    docx_content = docx_file.read()

                # Use synchronous method to extract text from DOCX
                text = self._extract_from_docx_sync(docx_content)

                return text

            finally:
                # Clean up temporary files
                try:
                    os.unlink(temp_doc_path)
                except Exception as e:
                    logger.warning(
                        f"Failed to remove temporary file {temp_doc_path}: {e}"
                    )

                try:
                    import shutil

                    shutil.rmtree(temp_dir, ignore_errors=True)
                except Exception as e:
                    logger.warning(
                        f"Failed to remove temporary directory {temp_dir}: {e}"
                    )

        except subprocess.TimeoutExpired:
            logger.error("LibreOffice conversion timeout")
            raise ValueError("DOC conversion timeout")
        except MemoryError as e:
            logger.error(f"LibreOffice exceeded memory limit: {str(e)}")
            raise ValueError("DOC conversion failed: memory limit exceeded")
        except Exception as e:
            logger.error(f"Error processing DOC: {str(e)}")
            raise ValueError(f"Error processing DOC: {str(e)}")

    def _extract_from_excel_sync(self, content: bytes) -> str:
        """Synchronous data extraction from Excel files."""
        if not pd:
            raise ImportError("pandas is not installed")

        try:
            excel_data = pd.read_excel(io.BytesIO(content), sheet_name=None)
            text_parts = []

            for sheet_name, df in excel_data.items():
                text_parts.append(f"[Sheet: {sheet_name}]")
                text_parts.append(df.to_csv(index=False))

            return "\n\n".join(text_parts)

        except Exception as e:
            logger.error(f"Error processing Excel: {str(e)}")
            raise ValueError(f"Error processing Excel: {str(e)}")

    def _extract_from_csv_sync(self, content: bytes) -> str:
        """Synchronous data extraction from CSV files."""
        if not pd:
            raise ImportError("pandas is not installed")

        try:
            df = pd.read_csv(io.BytesIO(content))
            return df.to_csv(index=False)

        except Exception as e:
            logger.error(f"Error processing CSV: {str(e)}")
            raise ValueError(f"Error processing CSV: {str(e)}")

    def _extract_from_pptx_sync(self, content: bytes) -> str:
        """Synchronous text extraction from PPTX with full extraction according to project specs."""
        if not Presentation:
            raise ImportError("python-pptx is not installed")

        try:
            prs = Presentation(io.BytesIO(content))
            text_parts = []

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                slide_text.append(f"[Slide {slide_num}]")

                # Extract text from slide shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)

                # Extract speaker notes - according to project specs
                try:
                    if hasattr(slide, "notes_slide") and slide.notes_slide:
                        notes_text = []
                        # Extract notes from text shapes
                        for shape in slide.notes_slide.shapes:
                            if hasattr(shape, "text") and shape.text.strip():
                                # Filter standard PowerPoint headers
                                if shape.text.strip() not in ["Заметки", "Notes"]:
                                    notes_text.append(shape.text.strip())

                        if notes_text:
                            slide_text.append(
                                f"[Speaker Notes]\n{' '.join(notes_text)}"
                            )
                except Exception as e:
                    logger.debug(
                        f"Failed to extract speaker notes from slide {slide_num}: {str(e)}"
                    )

                if len(slide_text) > 1:  # More than just slide title
                    text_parts.append("\n".join(slide_text))

            return "\n\n".join(text_parts)

        except Exception as e:
            logger.error(f"Error processing PPTX: {str(e)}")
            raise ValueError(f"Error processing PPTX: {str(e)}")

    def _extract_from_ppt_sync(self, content: bytes) -> str:
        """Synchronous text extraction from PPT via conversion to PPTX using LibreOffice."""
        if not Presentation:
            raise ImportError("python-pptx is not installed")

        try:
            # Create temporary files
            with tempfile.NamedTemporaryFile(suffix=".ppt", delete=False) as temp_ppt:
                temp_ppt.write(content)
                temp_ppt_path = temp_ppt.name

            # Create temporary output directory
            temp_dir = tempfile.mkdtemp()

            try:
                # Convert .ppt to .pptx using LibreOffice with resource limits
                from .config import settings
                from .utils import run_subprocess_with_limits

                result = run_subprocess_with_limits(
                    command=[
                        "libreoffice",
                        "--headless",
                        "--convert-to",
                        "pptx",
                        "--outdir",
                        temp_dir,
                        temp_ppt_path,
                    ],
                    timeout=30,
                    memory_limit=settings.MAX_LIBREOFFICE_MEMORY,
                    capture_output=True,
                    text=True,
                )

                if result.returncode != 0:
                    logger.error(f"LibreOffice conversion failed: {result.stderr}")
                    raise ValueError("Failed to convert PPT to PPTX")

                # Find converted file
                ppt_filename = os.path.splitext(os.path.basename(temp_ppt_path))[0]
                pptx_path = os.path.join(temp_dir, f"{ppt_filename}.pptx")

                if not os.path.exists(pptx_path):
                    raise ValueError("Converted PPTX file not found")

                # Read converted PPTX file
                with open(pptx_path, "rb") as pptx_file:
                    pptx_content = pptx_file.read()

                # Use synchronous method to extract text from PPTX
                text = self._extract_from_pptx_sync(pptx_content)

                return text

            finally:
                # Clean up temporary files
                try:
                    os.unlink(temp_ppt_path)
                except Exception as e:
                    logger.warning(
                        f"Failed to remove temporary file {temp_ppt_path}: {e}"
                    )

                try:
                    import shutil

                    shutil.rmtree(temp_dir, ignore_errors=True)
                except Exception as e:
                    logger.warning(
                        f"Failed to remove temporary directory {temp_dir}: {e}"
                    )

        except subprocess.TimeoutExpired:
            logger.error("LibreOffice conversion timeout")
            raise ValueError("PPT conversion timeout")
        except MemoryError as e:
            logger.error(f"LibreOffice exceeded memory limit: {str(e)}")
            raise ValueError("PPT conversion failed: memory limit exceeded")
        except Exception as e:
            logger.error(f"Error processing PPT: {str(e)}")
            raise ValueError(f"Error processing PPT: {str(e)}")

    def _extract_from_txt_sync(self, content: bytes) -> str:
        """Synchronous text extraction from TXT files."""
        try:
            return self._decode_text_content(content)
        except Exception as e:
            logger.error(f"Error processing TXT: {str(e)}")
            raise ValueError(f"Error processing TXT: {str(e)}")

    def _decode_text_content(self, content: bytes) -> str:
        """Decoding content with encoding auto-detection."""
        encodings = self._get_encoding_list()

        for encoding in encodings:
            decoded_text = self._try_decode_with_encoding(content, encoding)
            if decoded_text is not None:
                return decoded_text

        # If no encoding worked, use character replacement
        logger.warning(
            "Could not determine file encoding, using UTF-8 with character replacement"
        )
        return content.decode("utf-8", errors="replace")

    def _get_encoding_list(self) -> list:
        """Get list of encodings to check."""
        return [
            "utf-8",  # Standard UTF-8
            "mac-cyrillic",  # Macintosh encoding for Cyrillic
            "cp1251",  # Windows-1251 (main Windows encoding for Russian)
            "windows-1251",  # Alternative name for cp1251
            "koi8-r",  # KOI8-R (old Soviet encoding)
            "cp866",  # DOS encoding for Russian
            "iso-8859-5",  # ISO encoding for Cyrillic
            "utf-16",  # UTF-16 (sometimes used in Windows)
            "utf-16le",  # UTF-16 Little Endian
            "utf-16be",  # UTF-16 Big Endian
            "latin-1",  # ISO-8859-1 (fallback)
            "ascii",  # ASCII (basic encoding)
        ]

    def _try_decode_with_encoding(self, content: bytes, encoding: str) -> str:
        """Attempt to decode with quality check."""
        try:
            decoded_text = content.decode(encoding)

            if not self._is_decoding_quality_good(decoded_text):
                return None

            if not self._is_mac_cyrillic_valid(decoded_text, encoding):
                return None

            return decoded_text
        except UnicodeError:
            return None

    def _is_decoding_quality_good(self, text: str) -> bool:
        """Check decoding quality by replacement character count."""
        if "�" in text:
            replacement_ratio = text.count("�") / len(text)
            return replacement_ratio <= 0.1  # No more than 10% replacement characters
        return True

    def _is_mac_cyrillic_valid(self, text: str, encoding: str) -> bool:
        """Additional validation for mac-cyrillic encoding."""
        if encoding != "mac-cyrillic" or not text:
            return True

        if self._has_suspicious_start_chars(text):
            return False

        return self._has_valid_cyrillic_ratio(text)

    def _has_suspicious_start_chars(self, text: str) -> bool:
        """Check for suspicious characters at the start of text."""
        suspicious_chars = [
            '"',
            "'",
            "`",
            "«",
            "»",
            '"',
            '"',
            """, """,
            chr(8220),
            chr(8221),
        ]
        return len(text) > 1 and text[0] in suspicious_chars

    def _has_valid_cyrillic_ratio(self, text: str) -> bool:
        """Check the ratio of Cyrillic to Latin characters."""
        cyrillic_count = sum(1 for char in text if "\u0400" <= char <= "\u04ff")
        latin_count = sum(1 for char in text if "a" <= char.lower() <= "z")
        total_letters = cyrillic_count + latin_count

        if total_letters == 0:
            return True

        # If Cyrillic is less than 70% and present, it's suspicious
        return not (cyrillic_count / total_letters < 0.7 and cyrillic_count > 0)

    def _extract_from_source_code_sync(
        self, content: bytes, extension: str, filename: str
    ) -> str:
        """Synchronous text extraction from source code files."""
        try:
            # Decode file content
            text = self._decode_text_content(content)

            # Get programming language info and format result
            return self._format_source_code_output(text, extension, filename)

        except Exception as e:
            logger.error(f"Error processing source code {filename}: {str(e)}")
            raise ValueError(f"Error processing source code: {str(e)}")

    def _format_source_code_output(
        self, text: str, extension: str, filename: str
    ) -> str:
        """Formatting output for source code files."""
        language = self._get_programming_language(extension)
        header = self._create_source_code_header(language, filename, text)
        return header + "=" * 50 + "\n" + text

    def _get_programming_language(self, extension: str) -> str:
        """Determining programming language by file extension."""
        language_map = self._get_language_map()
        return language_map.get(extension.lower(), "Source Code")

    def _get_language_map(self) -> dict:
        """Get dictionary mapping extensions to programming languages."""
        return {
            # Python
            "py": "Python",
            "pyx": "Python",
            "pyi": "Python",
            "pyw": "Python",
            # JavaScript/TypeScript
            "js": "JavaScript",
            "jsx": "JavaScript",
            "ts": "TypeScript",
            "tsx": "TypeScript",
            "mjs": "JavaScript",
            "cjs": "JavaScript",
            # Java
            "java": "Java",
            "jav": "Java",
            # C/C++
            "c": "C",
            "cpp": "C++",
            "cxx": "C++",
            "cc": "C++",
            "c++": "C++",
            "h": "C Header",
            "hpp": "C++ Header",
            "hxx": "C++ Header",
            "h++": "C++ Header",
            # C#
            "cs": "C#",
            "csx": "C#",
            # PHP
            "php": "PHP",
            "php3": "PHP",
            "php4": "PHP",
            "php5": "PHP",
            "phtml": "PHP",
            # Ruby
            "rb": "Ruby",
            "rbw": "Ruby",
            "rake": "Ruby",
            "gemspec": "Ruby",
            # Go
            "go": "Go",
            "mod": "Go Module",
            "sum": "Go Sum",
            # Rust
            "rs": "Rust",
            "rlib": "Rust Library",
            # Swift
            "swift": "Swift",
            # Kotlin
            "kt": "Kotlin",
            "kts": "Kotlin Script",
            # Scala
            "scala": "Scala",
            "sc": "Scala",
            # R
            "r": "R",
            "R": "R",
            "rmd": "R Markdown",
            "Rmd": "R Markdown",
            # SQL
            "sql": "SQL",
            "ddl": "SQL DDL",
            "dml": "SQL DML",
            # Shell
            "sh": "Shell",
            "bash": "Bash",
            "zsh": "Zsh",
            "fish": "Fish",
            "ksh": "Ksh",
            "csh": "Csh",
            "tcsh": "Tcsh",
            # PowerShell
            "ps1": "PowerShell",
            "psm1": "PowerShell Module",
            "psd1": "PowerShell Data",
            # Perl
            "pl": "Perl",
            "pm": "Perl Module",
            "pod": "Perl Documentation",
            "t": "Perl Test",
            # Lua
            "lua": "Lua",
            # 1C and OneScript
            "bsl": "1C:Enterprise",
            "os": "OneScript",
            # Configuration files
            "ini": "INI Config",
            "cfg": "Config",
            "conf": "Config",
            "config": "Config",
            "toml": "TOML",
            "properties": "Properties",
            # Web
            "css": "CSS",
            "scss": "SCSS",
            "sass": "Sass",
            "less": "Less",
            "styl": "Stylus",
            # Markup
            "tex": "LaTeX",
            "latex": "LaTeX",
            "rst": "reStructuredText",
            "adoc": "AsciiDoc",
            "asciidoc": "AsciiDoc",
            # Data
            "jsonl": "JSON Lines",
            "ndjson": "NDJSON",
            "jsonc": "JSON with Comments",
            # Docker
            "dockerfile": "Dockerfile",
            "containerfile": "Containerfile",
            # Makefile
            "makefile": "Makefile",
            "mk": "Makefile",
            "mak": "Makefile",
            # Git
            "gitignore": "Git Ignore",
            "gitattributes": "Git Attributes",
            "gitmodules": "Git Modules",
        }

    def _create_source_code_header(
        self, language: str, filename: str, text: str
    ) -> str:
        """Creating a header for source code file."""
        header = f"=== {language} File: {filename} ===\n"

        lines = text.split("\n")
        line_count = len(lines)
        header += f"Lines: {line_count}\n"

        # If file is too long, add a warning
        if line_count > 1000:
            header += f"Warning: Large file with {line_count} lines\n"

        return header

    def _extract_from_html_sync(self, content: bytes) -> str:
        """Synchronous text extraction from HTML."""
        if not BeautifulSoup:
            raise ImportError("beautifulsoup4 is not installed")

        try:
            text = content.decode("utf-8", errors="replace")
            soup = BeautifulSoup(text, "html.parser")

            # Remove script and style tags
            for script in soup(["script", "style"]):
                script.decompose()

            # Get text
            text = soup.get_text()

            # Clean up extra spaces
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            return "\n".join(chunk for chunk in chunks if chunk)

        except Exception as e:
            logger.error(f"Error processing HTML: {str(e)}")
            raise ValueError(f"Error processing HTML: {str(e)}")

    def _extract_from_markdown_sync(self, content: bytes) -> str:
        """Synchronous text extraction from Markdown."""
        try:
            text = content.decode("utf-8", errors="replace")

            if markdown:
                # Convert to HTML and extract text
                html = markdown.markdown(text)
                if BeautifulSoup:
                    soup = BeautifulSoup(html, "html.parser")
                    return soup.get_text()

            # If markdown is not installed, return as is
            return text

        except Exception as e:
            logger.error(f"Error processing Markdown: {str(e)}")
            raise ValueError(f"Error processing Markdown: {str(e)}")

    def _extract_from_json_sync(self, content: bytes) -> str:
        """Synchronous text extraction from JSON."""
        import json

        try:
            text = content.decode("utf-8", errors="replace")
            data = json.loads(text)

            # Recursive extraction of all string values
            def extract_strings(obj, path=""):
                strings = []
                if isinstance(obj, dict):
                    for key, value in obj.items():
                        new_path = f"{path}.{key}" if path else key
                        strings.extend(extract_strings(value, new_path))
                elif isinstance(obj, list):
                    for i, value in enumerate(obj):
                        new_path = f"{path}[{i}]" if path else f"[{i}]"
                        strings.extend(extract_strings(value, new_path))
                elif isinstance(obj, str):
                    if obj.strip():
                        strings.append(f"{path}: {obj}")
                return strings

            strings = extract_strings(data)
            return "\n".join(strings)

        except Exception as e:
            logger.error(f"Error processing JSON: {str(e)}")
            raise ValueError(f"Error processing JSON: {str(e)}")

    def _extract_from_rtf_sync(self, content: bytes) -> str:
        """Synchronous text extraction from RTF."""
        if not rtf_to_text:
            raise ImportError("striprtf is not installed")

        try:
            text = content.decode("utf-8", errors="replace")
            plain_text = rtf_to_text(text)
            return plain_text

        except Exception as e:
            logger.error(f"Error processing RTF: {str(e)}")
            raise ValueError(f"Error processing RTF: {str(e)}")

    def _extract_from_xml_sync(self, content: bytes) -> str:
        """Synchronous text extraction from XML."""
        try:
            text = content.decode("utf-8", errors="replace")
            root = ET.fromstring(text)

            # Recursive extraction of all elements and attributes
            def extract_from_element(elem, path=""):
                strings = []

                current_path = f"{path}.{elem.tag}" if path else elem.tag

                # Add element text
                if elem.text and elem.text.strip():
                    strings.append(f"{current_path}: {elem.text.strip()}")

                # Add attributes
                for attr_name, attr_value in elem.attrib.items():
                    if attr_value.strip():
                        strings.append(f"{current_path}@{attr_name}: {attr_value}")

                # Recursively process child elements
                for child in elem:
                    strings.extend(extract_from_element(child, current_path))

                return strings

            strings = extract_from_element(root)
            return "\n".join(strings)

        except Exception as e:
            logger.error(f"Error processing XML: {str(e)}")
            raise ValueError(f"Error processing XML: {str(e)}")

    def _extract_from_yaml_sync(self, content: bytes) -> str:
        """Synchronous text extraction from YAML."""
        if not yaml:
            raise ImportError("PyYAML is not installed")

        try:
            text = content.decode("utf-8", errors="replace")
            data = yaml.safe_load(text)
            strings = self._extract_yaml_strings(data)
            return "\n".join(strings)

        except Exception as e:
            logger.error(f"Error processing YAML: {str(e)}")
            raise ValueError(f"Error processing YAML: {str(e)}")

    def _extract_yaml_strings(self, obj, path="") -> list:
        """Recursive extraction of all string values from YAML."""
        strings = []

        if isinstance(obj, dict):
            strings.extend(self._extract_yaml_dict_strings(obj, path))
        elif isinstance(obj, list):
            strings.extend(self._extract_yaml_list_strings(obj, path))
        elif isinstance(obj, str) and obj.strip():
            strings.append(f"{path}: {obj}")

        return strings

    def _extract_yaml_dict_strings(self, obj_dict: dict, path: str) -> list:
        """Extract strings from YAML dictionary."""
        strings = []
        for key, value in obj_dict.items():
            new_path = f"{path}.{key}" if path else key
            strings.extend(self._extract_yaml_strings(value, new_path))
        return strings

    def _extract_yaml_list_strings(self, obj_list: list, path: str) -> list:
        """Extract strings from YAML list."""
        strings = []
        for i, value in enumerate(obj_list):
            new_path = f"{path}[{i}]" if path else f"[{i}]"
            strings.extend(self._extract_yaml_strings(value, new_path))
        return strings

    def _extract_from_odt_sync(self, content: bytes) -> str:
        """Synchronous text extraction from ODT."""
        if not load:
            raise ImportError("odfpy is not installed")

        temp_file_path = None
        try:
            with tempfile.NamedTemporaryFile(suffix=".odt", delete=False) as temp_file:
                temp_file.write(content)
                temp_file_path = temp_file.name

            doc = load(temp_file_path)
            text_parts = []

            # Extract all text elements
            for p in doc.getElementsByType(P):
                text = extractText(p)
                if text.strip():
                    text_parts.append(text)

            return "\n".join(text_parts)

        except Exception as e:
            logger.error(f"Error processing ODT: {str(e)}")
            raise ValueError(f"Error processing ODT: {str(e)}")
        finally:
            # Guaranteed removal of temporary file
            if temp_file_path and os.path.exists(temp_file_path):
                try:
                    os.unlink(temp_file_path)
                except OSError as e:
                    logger.warning(
                        f"Failed to remove temporary file {temp_file_path}: {str(e)}"
                    )

    def _extract_from_epub_sync(self, content: bytes) -> str:
        """Synchronous text extraction from EPUB."""
        if not BeautifulSoup:
            raise ImportError("beautifulsoup4 is not installed")

        try:
            text_parts = []
            extracted_size = 0

            with zipfile.ZipFile(io.BytesIO(content), "r") as zip_ref:
                for file_info in zip_ref.infolist():
                    if self._should_stop_epub_extraction(
                        extracted_size, file_info.file_size
                    ):
                        break

                    if self._is_epub_html_file(file_info.filename):
                        text, new_size = self._extract_epub_html_text(
                            zip_ref, file_info
                        )
                        if text:
                            text_parts.append(text)
                        extracted_size += new_size

            return "\n\n".join(text_parts)

        except Exception as e:
            logger.error(f"Error processing EPUB: {str(e)}")
            raise ValueError(f"Error processing EPUB: {str(e)}")

    def _should_stop_epub_extraction(self, current_size: int, file_size: int) -> bool:
        """Check size limit for EPUB."""
        if current_size + file_size > settings.MAX_EXTRACTED_SIZE:
            logger.warning("EPUB uncompressed content size limit reached")
            return True
        return False

    def _is_epub_html_file(self, filename: str) -> bool:
        """Check if file is HTML for EPUB."""
        return filename.endswith((".html", ".xhtml", ".htm"))

    def _extract_epub_html_text(self, zip_ref, file_info) -> tuple:
        """Extract text from EPUB HTML file."""
        try:
            html_content = zip_ref.read(file_info.filename)
            html_text = html_content.decode("utf-8", errors="replace")

            # HTML parsing
            soup = BeautifulSoup(html_text, "html.parser")

            # Remove script and style tags
            for script in soup(["script", "style"]):
                script.decompose()

            # Extract text
            text = soup.get_text()
            return text.strip() if text.strip() else None, file_info.file_size

        except Exception as e:
            logger.warning(f"Error processing file {file_info.filename}: {e}")
            return None, 0

    def _extract_from_eml_sync(self, content: bytes) -> str:
        """Synchronous text extraction from EML."""
        import email

        try:
            msg_text = self._decode_eml_content(content)
            msg = email.message_from_string(msg_text)
            text_parts = []

            # Extract headers
            text_parts.extend(self._extract_eml_headers(msg))
            text_parts.append("---")

            # Extract email body
            if msg.is_multipart():
                text_parts.extend(self._extract_eml_body_multipart(msg))
            else:
                text_parts.extend(self._extract_eml_body_simple(msg))

            return (
                "\n".join(text_parts)
                if text_parts
                else "Could not extract readable text from EML file"
            )

        except Exception as e:
            logger.error(f"Error processing EML: {str(e)}")
            raise ValueError(f"Error processing EML: {str(e)}")

    def _decode_eml_content(self, content: bytes) -> str:
        """Decoding EML file content."""
        for encoding in ["utf-8", "cp1251", "latin-1"]:
            try:
                return content.decode(encoding)
            except UnicodeDecodeError:
                continue
        return content.decode("utf-8", errors="replace")

    def _extract_eml_headers(self, msg) -> list:
        """Extract headers from EML."""
        from email.header import decode_header

        text_parts = []
        headers = ["From", "To", "Subject", "Date"]

        for header in headers:
            value = msg.get(header)
            if value:
                decoded_value = self._decode_eml_header(value, decode_header)
                text_parts.append(f"{header}: {decoded_value}")

        return text_parts

    def _decode_eml_header(self, value: str, decode_header_func) -> str:
        """Decoding EML header."""
        decoded_parts = decode_header_func(value)
        decoded_value = ""

        for part, encoding in decoded_parts:
            if isinstance(part, bytes):
                if encoding:
                    decoded_value += part.decode(encoding)
                else:
                    decoded_value += part.decode("utf-8", errors="replace")
            else:
                decoded_value += part

        return decoded_value

    def _extract_eml_body_multipart(self, msg) -> list:
        """Extract body of multipart EML message."""
        text_parts = []

        for part in msg.walk():
            content_type = part.get_content_type()
            if content_type in ["text/plain", "text/html"]:
                try:
                    body_text = self._extract_eml_part_text(part, content_type)
                    if body_text and body_text.strip():
                        text_parts.append(body_text)
                except Exception as e:
                    logger.warning(f"Error processing message part: {e}")

        return text_parts

    def _extract_eml_body_simple(self, msg) -> list:
        """Extract body of simple EML message."""
        text_parts = []

        try:
            payload = msg.get_payload(decode=True)
            if payload:
                charset = msg.get_content_charset() or "utf-8"
                body_text = self._decode_payload(payload, charset)
                if body_text.strip():
                    text_parts.append(body_text)
        except Exception as e:
            logger.warning(f"Error processing message body: {e}")

        return text_parts

    def _extract_eml_part_text(self, part, content_type: str) -> str:
        """Extract text from EML part."""
        payload = part.get_payload(decode=True)
        if not payload:
            return ""

        charset = part.get_content_charset() or "utf-8"
        body_text = self._decode_payload(payload, charset)

        # HTML processing
        if content_type == "text/html" and BeautifulSoup:
            soup = BeautifulSoup(body_text, "html.parser")
            body_text = soup.get_text()

        return body_text

    def _decode_payload(self, payload: bytes, charset: str) -> str:
        """Decoding payload with error handling."""
        try:
            return payload.decode(charset)
        except UnicodeDecodeError:
            return payload.decode("utf-8", errors="replace")

    def _extract_from_msg_sync(self, content: bytes) -> str:
        """Synchronous text extraction from MSG."""
        try:
            text_parts = []

            # Extract UTF-16 text
            text_parts.extend(self._extract_utf16_text_from_msg(content))

            # Alternative approach - searching for ASCII text
            text_parts.extend(self._extract_ascii_text_from_msg(content, text_parts))

            return (
                "\n".join(text_parts)
                if text_parts
                else "Could not extract readable text from MSG file"
            )

        except Exception as e:
            logger.error(f"Error processing MSG: {str(e)}")
            raise ValueError(f"Error processing MSG: {str(e)}")

    def _extract_utf16_text_from_msg(self, content: bytes) -> list:
        """Extract UTF-16 text from MSG file."""
        text_parts = []
        try:
            text = content.decode("utf-16le", errors="ignore")
            lines = text.split("\n")
            clean_lines = self._clean_msg_lines(lines)
            unique_lines = self._filter_unique_lines(clean_lines, min_length=5)
            text_parts.extend(unique_lines)
        except Exception as e:
            logger.warning(f"Error decoding UTF-16: {e}")
        return text_parts

    def _clean_msg_lines(self, lines: list) -> list:
        """Cleaning MSG lines from control characters."""
        clean_lines = []
        for line in lines:
            # Remove null bytes and control characters
            clean_line = "".join(
                char for char in line if ord(char) >= 32 or char in "\t\n\r"
            )
            clean_line = clean_line.strip()

            # Skip lines that are too short or meaningless
            if self._is_valid_msg_line(clean_line):
                clean_lines.append(clean_line)

        return clean_lines

    def _is_valid_msg_line(self, line: str) -> bool:
        """Checking MSG line validity."""
        return (
            len(line) > 3
            and not line.startswith(("_", "\x00"))
            and any(c.isalpha() for c in line)
        )

    def _filter_unique_lines(self, lines: list, min_length: int = 5) -> list:
        """Filtering unique lines with a minimum length."""
        unique_lines = []
        seen = set()
        for line in lines:
            if line not in seen and len(line) > min_length:
                unique_lines.append(line)
                seen.add(line)
        return unique_lines

    def _extract_ascii_text_from_msg(
        self, content: bytes, existing_text_parts: list
    ) -> list:
        """Extract ASCII text from MSG file."""
        text_parts = []
        try:
            ascii_text = content.decode("ascii", errors="ignore")
            lines = ascii_text.split("\n")

            for line in lines:
                clean_line = line.strip()
                if self._is_valid_ascii_line(clean_line, existing_text_parts):
                    text_parts.append(clean_line)
        except Exception as e:
            logger.warning(f"Error extracting ASCII: {e}")
        return text_parts

    def _is_valid_ascii_line(self, line: str, existing_parts: list) -> bool:
        """Checking ASCII line validity."""
        return (
            len(line) > 10
            and any(c.isalpha() for c in line)
            and line not in existing_parts
        )

    def _safe_tesseract_ocr(self, image, temp_image_path: str = None) -> str:
        """
        Safely call Tesseract with resource limits.

        Args:
            image: PIL Image object
            temp_image_path: Path to temporary file (if None, created automatically)

        Returns:
            str: Recognized text
        """
        import os
        import tempfile

        from .config import settings
        from .utils import run_subprocess_with_limits

        temp_file_created = False

        try:
            # If temporary image path is not provided, create it
            if temp_image_path is None:
                with tempfile.NamedTemporaryFile(
                    suffix=".png", delete=False
                ) as temp_file:
                    temp_image_path = temp_file.name
                    temp_file_created = True

                # Save image to temporary file
                # Convert to RGB for PNG compatibility
                if image.mode in ("RGBA", "LA", "P"):
                    image = image.convert("RGB")
                image.save(temp_image_path, "PNG")

            # Create temporary output file
            with tempfile.NamedTemporaryFile(
                suffix=".txt", delete=False
            ) as output_file:
                output_path = output_file.name

            try:
                # Call Tesseract via safe function
                result = run_subprocess_with_limits(
                    command=[
                        "tesseract",
                        temp_image_path,
                        output_path.replace(".txt", ""),
                        "-l",
                        self.ocr_languages,
                    ],
                    timeout=30,
                    memory_limit=settings.MAX_TESSERACT_MEMORY,
                    capture_output=True,
                    text=True,
                )

                if result.returncode != 0:
                    logger.warning(
                        f"Tesseract exited with code {result.returncode}: {result.stderr}"
                    )
                    return ""

                # Read OCR result
                if os.path.exists(output_path):
                    with open(output_path, "r", encoding="utf-8") as f:
                        return f.read().strip()
                else:
                    logger.warning("OCR result file not found")
                    return ""

            finally:
                # Remove temporary output file
                try:
                    if os.path.exists(output_path):
                        os.unlink(output_path)
                except Exception as e:
                    logger.warning(
                        f"Failed to remove temporary file {output_path}: {e}"
                    )

        except subprocess.TimeoutExpired:
            logger.error("Tesseract OCR timeout")
            return ""
        except MemoryError as e:
            logger.error(f"Tesseract exceeded memory limit: {str(e)}")
            return ""
        except Exception as e:
            logger.error(f"OCR error: {str(e)}")
            return ""
        finally:
            # Remove temporary image file if we created it
            if (
                temp_file_created
                and temp_image_path
                and os.path.exists(temp_image_path)
            ):
                try:
                    os.unlink(temp_image_path)
                except Exception as e:
                    logger.warning(
                        f"Failed to remove temporary file {temp_image_path}: {e}"
                    )

    def _extract_from_image_sync(self, content: bytes) -> str:
        """Synchronous image OCR."""
        if not Image:
            raise ImportError("PIL is not installed")

        image = None
        try:
            # Image validation to prevent DoS attacks
            from .utils import validate_image_for_ocr

            is_valid, error_message = validate_image_for_ocr(content)
            if not is_valid:
                logger.warning(f"Image validation failed: {error_message}")
                raise ValueError(f"Image validation failed: {error_message}")

            image = Image.open(io.BytesIO(content))

            # Safe OCR with resource limits
            text = self._safe_tesseract_ocr(image)
            return text

        except Exception as e:
            logger.error(f"Error during image OCR: {str(e)}")
            raise ValueError(f"Error processing image: {str(e)}")
        finally:
            # CRITICAL FIX: explicitly close image to free memory
            if image:
                try:
                    image.close()
                except Exception as close_error:
                    logger.warning(f"Error closing image: {str(close_error)}")

    def _check_mime_type(self, content: bytes, filename: str) -> bool:
        """Check file MIME type to prevent extension spoofing."""
        import mimetypes

        try:
            # Determine MIME type by content (magic bytes)
            mime_signatures = {
                b"\x50\x4b\x03\x04": [
                    "application/zip",
                    "application/epub+zip",
                    "application/vnd.openxmlformats",
                ],
                b"\x50\x4b\x07\x08": ["application/zip", "application/epub+zip"],
                b"\x50\x4b\x05\x06": ["application/zip", "application/epub+zip"],
                b"%PDF": ["application/pdf"],
                b"\xd0\xcf\x11\xe0": [
                    "application/msword",
                    "application/vnd.ms-excel",
                    "application/vnd.ms-powerpoint",
                ],
                b"\x89PNG": ["image/png"],
                b"\xff\xd8\xff": ["image/jpeg"],
                b"GIF8": ["image/gif"],
                b"BM": ["image/bmp"],
                b"II*\x00": ["image/tiff"],
                b"MM\x00*": ["image/tiff"],
                b"<!DOCTYPE": ["text/html"],
                b"<html": ["text/html"],
                b"<?xml": ["text/xml", "application/xml"],
            }

            # Check file signature
            file_start = content[:10]
            detected_mime = None

            for signature, mime_types in mime_signatures.items():
                if file_start.startswith(signature):
                    detected_mime = mime_types[0]
                    break

            # Determine expected MIME type by extension
            expected_mime, _ = mimetypes.guess_type(filename)

            # If we cannot determine MIME type, allow it
            if not detected_mime or not expected_mime:
                return True

            # Check consistency
            return detected_mime in mime_signatures.get(file_start[:4], [expected_mime])

        except Exception as e:
            logger.warning(f"Error during MIME type check: {str(e)}")
            return True  # Allow processing in case of error

    def _extract_from_archive(
        self, content: bytes, filename: str, nesting_level: int = 0
    ) -> List[Dict[str, Any]]:
        """Safe extraction of files from archive."""
        # Check nesting depth
        if nesting_level >= settings.MAX_ARCHIVE_NESTING:
            logger.warning(
                f"Maximum archive nesting depth exceeded: {filename}"
            )
            raise ValueError("Maximum archive nesting level exceeded")

        # Check archive size
        if len(content) > settings.MAX_ARCHIVE_SIZE:
            logger.warning(f"Archive {filename} too large: {len(content)} bytes")
            raise ValueError("Archive size exceeds maximum allowed size")

        extension = get_file_extension(filename)
        logger.info(
            f"Processing archive {filename} (type: {extension}, size: {len(content)} bytes)"
        )

        extracted_files = []

        # Create temporary directory for safe operation
        with tempfile.TemporaryDirectory() as temp_dir:
            temp_path = Path(temp_dir)
            archive_path = temp_path / f"archive_{int(time.time())}.{extension}"

            try:
                # Write archive to temporary file
                with open(archive_path, "wb") as f:
                    f.write(content)

                # Extract files depending on archive type
                extract_dir = temp_path / "extracted"
                extract_dir.mkdir(exist_ok=True)

                if extension == "zip":
                    extracted_files = self._extract_zip_files(
                        archive_path, extract_dir, filename, nesting_level
                    )
                elif extension in [
                    "tar",
                    "gz",
                    "bz2",
                    "xz",
                    "tar.gz",
                    "tar.bz2",
                    "tar.xz",
                    "tgz",
                    "tbz2",
                    "txz",
                ]:
                    extracted_files = self._extract_tar_files(
                        archive_path, extract_dir, filename, nesting_level
                    )
                elif extension == "rar":
                    extracted_files = self._extract_rar_files(
                        archive_path, extract_dir, filename, nesting_level
                    )
                elif extension == "7z":
                    extracted_files = self._extract_7z_files(
                        archive_path, extract_dir, filename, nesting_level
                    )
                else:
                    raise ValueError(f"Unsupported archive format: {extension}")

                logger.info(
                    f"Successfully processed {len(extracted_files)} files from archive {filename}"
                )
                return extracted_files

            except Exception as e:
                logger.error(f"Error processing archive {filename}: {str(e)}")
                raise ValueError(f"Error processing archive: {str(e)}")

    def _extract_zip_files(
        self,
        archive_path: Path,
        extract_dir: Path,
        archive_name: str,
        nesting_level: int,
    ) -> List[Dict[str, Any]]:
        """Extract files from ZIP archive."""
        try:
            with zipfile.ZipFile(archive_path, "r") as zip_ref:
                self._validate_zip_size(zip_ref)
                return self._process_zip_files(
                    zip_ref, extract_dir, archive_name, nesting_level
                )
        except zipfile.BadZipFile:
            raise ValueError("Invalid ZIP file")

    def _validate_zip_size(self, zip_ref) -> None:
        """Check file sizes in ZIP archive for zip bomb protection."""
        total_size = 0
        for info in zip_ref.infolist():
            if not info.is_dir():
                total_size += info.file_size
                if total_size > settings.MAX_EXTRACTED_SIZE:
                    raise ValueError(
                        "Extracted files size exceeds maximum allowed size (zip bomb protection)"
                    )

    def _process_zip_files(
        self, zip_ref, extract_dir: Path, archive_name: str, nesting_level: int
    ) -> List[Dict[str, Any]]:
        """Process all files in ZIP archive."""
        extracted_files = []

        for info in zip_ref.infolist():
            if info.is_dir():
                continue

            file_result = self._extract_single_zip_file(
                info, zip_ref, extract_dir, archive_name, nesting_level
            )
            if file_result:
                extracted_files.extend(file_result)

        return extracted_files

    def _extract_single_zip_file(
        self, info, zip_ref, extract_dir: Path, archive_name: str, nesting_level: int
    ) -> List[Dict[str, Any]]:
        """Extract and process a single file from ZIP archive."""
        # Sanitize filename
        safe_filename = self._sanitize_archive_filename(info.filename)
        if not safe_filename:
            return []

        # Filter system files
        if self._is_system_file(safe_filename):
            return []

        # Create safe path for extraction
        safe_path = extract_dir / safe_filename
        safe_path.parent.mkdir(parents=True, exist_ok=True)

        try:
            # Extract file
            with zip_ref.open(info) as source, open(safe_path, "wb") as target:
                shutil.copyfileobj(source, target)

            # Process file
            file_content = safe_path.read_bytes()
            return (
                self._process_extracted_file(
                    file_content,
                    safe_filename,
                    safe_path.name,
                    archive_name,
                    nesting_level,
                )
                or []
            )

        except Exception as e:
            logger.warning(
                f"Error processing file {safe_filename} from archive {archive_name}: {str(e)}"
            )
            return []

    def _extract_tar_files(
        self,
        archive_path: Path,
        extract_dir: Path,
        archive_name: str,
        nesting_level: int,
    ) -> List[Dict[str, Any]]:
        """Extract files from TAR archive."""
        extracted_files = []
        total_size = 0

        try:
            with tarfile.open(archive_path, "r:*") as tar_ref:
                # Check uncompressed file sizes
                for member in tar_ref.getmembers():
                    if member.isfile():
                        total_size += member.size

                        if total_size > settings.MAX_EXTRACTED_SIZE:
                            raise ValueError(
                                "Extracted files size exceeds maximum allowed size (tar bomb protection)"
                            )

                # Extract files
                for member in tar_ref.getmembers():
                    if not member.isfile():
                        continue

                    # Sanitize filename
                    safe_filename = self._sanitize_archive_filename(member.name)
                    if not safe_filename:
                        continue

                    # Filter system files
                    if self._is_system_file(safe_filename):
                        continue

                    # Create safe path for extraction
                    safe_path = extract_dir / safe_filename
                    safe_path.parent.mkdir(parents=True, exist_ok=True)

                    try:
                        # Extract file
                        with (
                            tar_ref.extractfile(member) as source,
                            open(safe_path, "wb") as target,
                        ):
                            if source:
                                shutil.copyfileobj(source, target)

                        # Process file
                        file_content = safe_path.read_bytes()
                        file_result = self._process_extracted_file(
                            file_content,
                            safe_filename,
                            safe_path.name,
                            archive_name,
                            nesting_level,
                        )

                        if file_result:
                            extracted_files.extend(file_result)

                    except Exception as e:
                        logger.warning(
                            f"Error processing file {safe_filename} from archive {archive_name}: {str(e)}"
                        )
                        continue

        except tarfile.TarError:
            raise ValueError("Invalid TAR file")

        return extracted_files

    def _extract_rar_files(
        self,
        archive_path: Path,
        extract_dir: Path,
        archive_name: str,
        nesting_level: int,
    ) -> List[Dict[str, Any]]:
        """Extract files from RAR archive."""
        if not rarfile:
            raise ValueError("RAR support not available. Install rarfile library.")

        extracted_files = []
        total_size = 0

        try:
            with rarfile.RarFile(archive_path, "r") as rar_ref:
                # Check uncompressed file sizes
                for info in rar_ref.infolist():
                    if info.is_dir():
                        continue
                    total_size += info.file_size

                    if total_size > settings.MAX_EXTRACTED_SIZE:
                        raise ValueError(
                            "Extracted files size exceeds maximum allowed size (rar bomb protection)"
                        )

                # Extract files
                for info in rar_ref.infolist():
                    if info.is_dir():
                        continue

                    # Sanitize filename
                    safe_filename = self._sanitize_archive_filename(info.filename)
                    if not safe_filename:
                        continue

                    # Filter system files
                    if self._is_system_file(safe_filename):
                        continue

                    # Create safe path for extraction
                    safe_path = extract_dir / safe_filename
                    safe_path.parent.mkdir(parents=True, exist_ok=True)

                    try:
                        # Extract file
                        with (
                            rar_ref.open(info) as source,
                            open(safe_path, "wb") as target,
                        ):
                            shutil.copyfileobj(source, target)

                        # Process file
                        file_content = safe_path.read_bytes()
                        file_result = self._process_extracted_file(
                            file_content,
                            safe_filename,
                            safe_path.name,
                            archive_name,
                            nesting_level,
                        )

                        if file_result:
                            extracted_files.extend(file_result)

                    except Exception as e:
                        logger.warning(
                            f"Error processing file {safe_filename} from archive {archive_name}: {str(e)}"
                        )
                        continue

        except rarfile.RarError:
            raise ValueError("Invalid RAR file")

        return extracted_files

    def _extract_7z_files(
        self,
        archive_path: Path,
        extract_dir: Path,
        archive_name: str,
        nesting_level: int,
    ) -> List[Dict[str, Any]]:
        """Extract files from 7Z archive."""
        if not py7zr:
            raise ValueError("7Z support not available. Install py7zr library.")

        extracted_files = []
        total_size = 0

        try:
            with py7zr.SevenZipFile(archive_path, "r") as sz_ref:
                # Check uncompressed file sizes
                for info in sz_ref.list():
                    if info.is_dir:
                        continue
                    total_size += info.uncompressed

                    if total_size > settings.MAX_EXTRACTED_SIZE:
                        raise ValueError(
                            "Extracted files size exceeds maximum allowed size (7z bomb protection)"
                        )

                # Extract files
                sz_ref.extractall(extract_dir)

                # Process extracted files
                for root, _dirs, files in os.walk(extract_dir):
                    for file in files:
                        file_path = Path(root) / file
                        relative_path = file_path.relative_to(extract_dir)

                        # Sanitize filename
                        safe_filename = self._sanitize_archive_filename(
                            str(relative_path)
                        )
                        if not safe_filename:
                            continue

                        # Filter system files
                        if self._is_system_file(safe_filename):
                            continue

                        try:
                            # Process file
                            file_content = file_path.read_bytes()
                            file_result = self._process_extracted_file(
                                file_content,
                                safe_filename,
                                file_path.name,
                                archive_name,
                                nesting_level,
                            )

                            if file_result:
                                extracted_files.extend(file_result)

                        except Exception as e:
                            logger.warning(
                                f"Error processing file {safe_filename} from archive {archive_name}: {str(e)}"
                            )
                            continue

        except py7zr.Bad7zFile:
            raise ValueError("Invalid 7Z file")

        return extracted_files

    def _process_extracted_file(
        self,
        content: bytes,
        filename: str,
        basename: str,
        archive_name: str,
        nesting_level: int,
    ) -> Optional[List[Dict[str, Any]]]:
        """Process extracted file."""
        try:
            # If file is an archive, recursively process it
            if is_archive_format(basename, settings.SUPPORTED_FORMATS):
                return self._extract_from_archive(content, basename, nesting_level + 1)

            # If file is supported, extract text
            if is_supported_format(basename, settings.SUPPORTED_FORMATS):
                extension = get_file_extension(basename)
                text = self._extract_text_by_format(content, extension, basename)

                return [
                    {
                        "filename": basename,
                        "path": f"{archive_name}/{filename}",
                        "size": len(content),
                        "type": extension,
                        "text": text.strip() if text else "",
                    }
                ]

            return None

        except Exception as e:
            logger.warning(f"Error processing file {filename}: {str(e)}")
            return None

    def _sanitize_archive_filename(self, filename: str) -> str:
        """Sanitize archive filename."""
        if not filename:
            return ""

        # Remove dangerous paths
        filename = filename.replace("..", "").replace("\\", "/").strip("/")
        
        # Check for absolute paths
        if filename.startswith("/"):
            filename = filename[1:]

        # Remove empty path segments
        parts = [part for part in filename.split("/") if part and part != "."]

        if not parts:
            return ""

        return "/".join(parts)

    def _is_system_file(self, filename: str) -> bool:
        """Check if file is a system file."""
        system_files = [
            ".DS_Store",
            "Thumbs.db",
            ".git/",
            ".svn/",
            ".hg/",
            "__MACOSX/",
            ".localized",
            "desktop.ini",
            "folder.ini",
        ]

        filename_lower = filename.lower()
        for system_file in system_files:
            if system_file in filename_lower:
                return True

        return False

    def _ocr_from_pdf_image_sync(self, page, img_info) -> str:
        """Synchronous OCR of image from PDF."""
        if not Image:
            return ""

        try:
            # Get image coordinates
            x0, y0, x1, y1 = (
                img_info["x0"],
                img_info["y0"],
                img_info["x1"],
                img_info["y1"],
            )

            # Check if area size is reasonable
            width = abs(x1 - x0)
            height = abs(y1 - y0)

            # Limit area size to prevent DoS
            max_dimension = 5000  # maximum size along any axis
            if width > max_dimension or height > max_dimension:
                logger.warning(f"Image area too large: {width}x{height}")
                return ""

            # Crop image area from entire page
            cropped_bbox = (x0, y0, x1, y1)
            cropped_page = page.crop(cropped_bbox)

            # Convert cropped area to high-resolution image
            img_pil = cropped_page.to_image(resolution=300)

            # Safe OCR with resource limits
            text = self._safe_tesseract_ocr(img_pil.original)

            return text

        except Exception as e:
            logger.warning(f"OCR error: {str(e)}")
            # Alternative approach - render entire page and crop
            try:
                # Convert entire page to PIL image
                page_image = page.to_image(resolution=300)
                pil_image = page_image.original  # Get PIL image

                # Calculate pixel coordinates (considering resolution=300)
                scale = 300 / 72  # PDF is usually 72 DPI, we render at 300 DPI
                pixel_bbox = (
                    int(x0 * scale),
                    int(y0 * scale),
                    int(x1 * scale),
                    int(y1 * scale),
                )

                # Check if pixel dimensions are reasonable
                pixel_width = abs(pixel_bbox[2] - pixel_bbox[0])
                pixel_height = abs(pixel_bbox[3] - pixel_bbox[1])

                if pixel_width * pixel_height > 25000000:  # 25MP максимум
                    logger.warning(
                        f"Image area too large: {pixel_width}x{pixel_height} pixels"
                    )
                    return ""

                # Crop image area
                cropped_img = pil_image.crop(pixel_bbox)

                # Safe OCR with resource limits
                text = self._safe_tesseract_ocr(cropped_img)

                return text

            except Exception as e2:
                logger.warning(
                    f"Alternative OCR attempt also failed: {str(e2)}"
                )
                return ""

    # Web extraction (new in v1.10.0)

    def _extract_page_with_playwright(
        self,
        url: str,
        user_agent: Optional[str] = None,
        extraction_options: Optional[Any] = None,
    ) -> tuple[str, str]:
        """
        Extract HTML content of page using Playwright (with JS support, updated in v1.10.2).

        Args:
            url: Page URL
            user_agent: Custom User-Agent
            extraction_options: Extraction settings

        Returns:
            tuple[str, str]: (html_content, final_url)
        """
        if not sync_playwright:
            raise ValueError("Playwright is not installed")

        # Define settings considering provided parameters or defaults
        web_page_timeout = (
            extraction_options.web_page_timeout
            if extraction_options and extraction_options.web_page_timeout is not None
            else settings.WEB_PAGE_TIMEOUT
        )

        js_render_timeout = (
            extraction_options.js_render_timeout
            if extraction_options and extraction_options.js_render_timeout is not None
            else settings.JS_RENDER_TIMEOUT
        )

        web_page_delay = (
            extraction_options.web_page_delay
            if extraction_options and extraction_options.web_page_delay is not None
            else settings.WEB_PAGE_DELAY
        )

        enable_lazy_loading_wait = (
            extraction_options.enable_lazy_loading_wait
            if extraction_options
            and extraction_options.enable_lazy_loading_wait is not None
            else settings.ENABLE_LAZY_LOADING_WAIT
        )

        html_content = ""
        final_url = url

        with sync_playwright() as p:
            # Launch Chromium (installed in Dockerfile)
            browser = p.chromium.launch(
                headless=True,
                args=[
                    "--no-sandbox",
                    "--disable-dev-shm-usage",
                    "--disable-gpu",
                    "--disable-extensions",
                    "--disable-web-security",  # for bypassing CORS in local development
                ],
            )

            try:
                # Determine whether to enable JavaScript
                enable_javascript = (
                    extraction_options.enable_javascript
                    if extraction_options
                    and extraction_options.enable_javascript is not None
                    else settings.ENABLE_JAVASCRIPT
                )

                context = browser.new_context(
                    user_agent=user_agent or settings.DEFAULT_USER_AGENT,
                    viewport={"width": 1280, "height": 720},
                    java_script_enabled=enable_javascript,  # Correct JS disable setting
                )

                page = None
                try:
                    page = context.new_page()

                    # Set timeouts for DoS protection
                    page.set_default_timeout(web_page_timeout * 1000)  # in milliseconds

                    # Additional DoS protection: limit JavaScript execution time
                    page.set_default_navigation_timeout(web_page_timeout * 1000)

                    # Navigate to page
                    logger.info(
                        f"Loading page with Playwright: {url} (JS: {'enabled' if enable_javascript else 'disabled'})"
                    )
                    response = page.goto(url, wait_until="domcontentloaded")

                    if not response.ok:
                        raise ValueError(f"HTTP {response.status}: {response.status_text}")

                    final_url = page.url

                    # Wait for additional JS loading (if enabled)
                    if enable_javascript:
                        logger.info(f"Waiting for JS rendering ({js_render_timeout}s)...")

                        # Wait for network load with DoS protection
                        try:
                            # Limit execution time for protection against resource-intensive scripts
                            page.wait_for_load_state(
                                "networkidle",
                                timeout=min(
                                    js_render_timeout * 1000, 15000
                                ),  # не более 15 сек
                            )
                        except Exception as e:
                            logger.warning(
                                f"Network wait timeout (DoS protection): {str(e)}"
                            )

                        # Lazy loading processing with infinite loop protection
                        if enable_lazy_loading_wait:
                            self._safe_scroll_for_lazy_loading(page, extraction_options)

                        # Additional delay for JS completion
                        import time

                        time.sleep(web_page_delay)

                    # Get final HTML
                    html_content = page.content()
                    logger.info(f"HTML received, size: {len(html_content)} characters")

                finally:
                    # CRITICAL FIX: explicitly close page to free memory
                    if page:
                        try:
                            page.close()
                        except Exception as e:
                            logger.warning(f"Error closing page: {str(e)}")
                    # Explicitly close context
                    try:
                        context.close()
                    except Exception as e:
                        logger.warning(f"Error closing context: {str(e)}")

            finally:
                browser.close()

        return html_content, final_url

    def _safe_scroll_for_lazy_loading(
        self, page, extraction_options: Optional[Any] = None
    ) -> None:
        """
        Safe page scroll to activate lazy loading with infinite loop protection (updated in v1.10.2).

        Args:
            page: Playwright page object
            extraction_options: Extraction settings
        """
        try:
            logger.info("Starting safe scroll for lazy loading activation...")

            # Determine maximum number of scroll attempts
            max_scroll_attempts = (
                extraction_options.max_scroll_attempts
                if extraction_options
                and extraction_options.max_scroll_attempts is not None
                else settings.MAX_SCROLL_ATTEMPTS
            )

            # Get initial page height
            initial_height = page.evaluate("document.body.scrollHeight")
            logger.info(f"Initial page height: {initial_height}px")

            scroll_attempts = 0
            last_height = initial_height
            stable_count = 0  # Counter of stable measurements

            while scroll_attempts < max_scroll_attempts:
                scroll_attempts += 1
                logger.info(f"Scroll attempt {scroll_attempts}/{max_scroll_attempts}")

                # Smooth scroll to end of page
                page.evaluate("window.scrollTo(0, document.body.scrollHeight)")

                # Wait a short delay for content to load
                import time

                time.sleep(1)

                # Check new height
                new_height = page.evaluate("document.body.scrollHeight")
                logger.info(f"New page height: {new_height}px")

                # If height hasn't changed, increment stability counter
                if new_height == last_height:
                    stable_count += 1
                    logger.info(f"Height stable, counter: {stable_count}")
                    
                    # If height has been stable 2 times in a row - stop
                    if stable_count >= 2:
                        logger.info(
                            "Page height stabilized, completing scroll"
                        )
                        break
                else:
                    # Height changed, reset counter
                    stable_count = 0
                    last_height = new_height

                # Additional check: if page grew too much, stop
                if new_height > initial_height * 10:  # If page grew 10x
                    logger.warning(
                        "Page grew suspiciously large, possibly infinite scroll"
                    )
                    break

            # Return to top of page
            page.evaluate("window.scrollTo(0, 0)")
            logger.info("Scroll completed, returned to top of page")

        except Exception as e:
            logger.warning(f"Error during lazy loading scroll: {str(e)}")

    def _determine_content_type(
        self,
        url: str,
        user_agent: Optional[str] = None,
        extraction_options: Optional[Any] = None,
    ) -> tuple[str, str]:
        """
        Determine content type via HEAD request (new in v1.10.3).

        Args:
            url: URL to check
            user_agent: Custom User-Agent
            extraction_options: Extraction settings

        Returns:
            tuple[str, str]: (content_type, final_url) - content type and final URL after redirects
        """
        if not requests:
            raise ValueError(
                "requests library not available for content type determination"
            )

        # Define settings considering provided parameters or defaults
        head_timeout = (
            extraction_options.web_page_timeout
            if extraction_options and extraction_options.web_page_timeout is not None
            else settings.HEAD_REQUEST_TIMEOUT
        )

        follow_redirects = (
            extraction_options.follow_redirects
            if extraction_options and extraction_options.follow_redirects is not None
            else True
        )

        max_redirects = (
            extraction_options.max_redirects
            if extraction_options and extraction_options.max_redirects is not None
            else 5
        )

        # Set User-Agent and headers
        headers = {
            "User-Agent": user_agent or settings.DEFAULT_USER_AGENT,
            "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8",
            "Accept-Language": "ru,en;q=0.5",
            "Accept-Encoding": "gzip, deflate",
            "Connection": "keep-alive",
        }

        # Create session to follow redirects
        session = requests.Session()
        session.headers.update(headers)

        try:
            logger.info(f"Performing HEAD request to determine content type: {url}")

            response = session.head(
                url, timeout=head_timeout, allow_redirects=follow_redirects, stream=True
            )

            if follow_redirects and len(response.history) > max_redirects:
                logger.warning(
                    f"Maximum number of redirects exceeded ({max_redirects})"
                )

            response.raise_for_status()

            content_type = response.headers.get("content-type", "").lower()
            final_url = response.url

            logger.info(f"Content-Type determined: {content_type} for URL: {final_url}")

            return content_type, final_url

        except Exception as e:
            logger.warning(f"HEAD request error: {str(e)}, trying GET request")
            # Fallback: perform a GET request but only read headers
            try:
                response = session.get(
                    url,
                    timeout=head_timeout,
                    allow_redirects=follow_redirects,
                    stream=True,
                )
                response.raise_for_status()

                content_type = response.headers.get("content-type", "").lower()
                final_url = response.url

                # Close connection without reading body
                response.close()

                logger.info(
                    f"Content-Type determined via GET: {content_type} for URL: {final_url}"
                )
                return content_type, final_url

            except Exception as get_error:
                logger.error(f"Error determining content type: {str(get_error)}")
                raise ValueError(f"Unable to determine content type: {str(get_error)}")
        finally:
            session.close()

    def _is_html_content(self, content_type: str, url: str) -> bool:
        """
        Determine if content is an HTML page (new in v1.10.3).

        Args:
            content_type: MIME type from headers
            url: URL for extension analysis as fallback

        Returns:
            bool: True if it is an HTML page
        """
        # Priority: Content-Type
        if "text/html" in content_type or "application/xhtml" in content_type:
            return True

        # Check specific cases
        if "text/plain" in content_type:
            # For text/plain check URL extension
            from app.utils import get_file_extension

            extension = get_file_extension(url.split("?")[0])  # remove parameters
            return extension in ["html", "htm"]

        # If Content-Type is undefined or missing
        if not content_type or "application/octet-stream" in content_type:
            # Use URL extension as fallback
            from app.utils import get_file_extension

            extension = get_file_extension(url.split("?")[0])  # remove parameters
            return (
                extension in ["html", "htm"] or extension is None
            )  # None likely means dynamic page

        return False

    def _download_and_extract_file(
        self,
        url: str,
        user_agent: Optional[str] = None,
        extraction_options: Optional[Any] = None,
    ) -> List[Dict[str, Any]]:
        """
        Download file from URL and process it as a regular file (new in v1.10.3).

        Args:
            url: URL of file to download
            user_agent: Custom User-Agent
            extraction_options: Extraction settings

        Returns:
            List[Dict[str, Any]]: Result of text extraction like from /v1/extract/file
        """
        if not requests:
            raise ValueError("requests library not available for file download")

        # Define settings considering provided parameters or defaults
        download_timeout = (
            extraction_options.web_page_timeout
            if extraction_options and extraction_options.web_page_timeout is not None
            else settings.FILE_DOWNLOAD_TIMEOUT
        )

        follow_redirects = (
            extraction_options.follow_redirects
            if extraction_options and extraction_options.follow_redirects is not None
            else True
        )

        # Set User-Agent and headers
        headers = {
            "User-Agent": user_agent or settings.DEFAULT_USER_AGENT,
            "Accept": "*/*",
            "Connection": "keep-alive",
        }

        # Create session
        session = requests.Session()
        session.headers.update(headers)

        temp_file_path = None
        try:
            logger.info(f"Downloading file from URL: {url}")

            response = session.get(
                url,
                timeout=download_timeout,
                allow_redirects=follow_redirects,
                stream=True,
            )
            response.raise_for_status()

            # Check file size
            content_length = response.headers.get("content-length")
            if content_length and int(content_length) > settings.MAX_FILE_SIZE:
                raise ValueError(
                    f"File too large: {content_length} bytes (max {settings.MAX_FILE_SIZE} bytes)"
                )

            # Determine filename
            filename = self._extract_filename_from_response(response, url)

            # Create temporary file
            import tempfile

            suffix = (
                f".{get_file_extension(filename)}"
                if get_file_extension(filename)
                else ""
            )
            with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as temp_file:
                temp_file_path = temp_file.name

                # Download file in chunks with size check
                downloaded_size = 0
                for chunk in response.iter_content(chunk_size=8192):
                    if chunk:
                        downloaded_size += len(chunk)
                        if downloaded_size > settings.MAX_FILE_SIZE:
                            raise ValueError(
                                f"File too large: exceeded {settings.MAX_FILE_SIZE} bytes during download"
                            )
                        temp_file.write(chunk)

            logger.info(f"File downloaded ({downloaded_size} bytes): {filename}")

            # Read downloaded file and process it as a regular file
            with open(temp_file_path, "rb") as f:
                file_content = f.read()

            # Use existing text extraction logic
            return self.extract_text(file_content, filename)

        except Exception as e:
            logger.error(f"Error downloading file {url}: {str(e)}")
            raise ValueError(f"Error downloading file: {str(e)}")
        finally:
            # Cleanup of temporary file
            if temp_file_path and os.path.exists(temp_file_path):
                try:
                    os.unlink(temp_file_path)
                    logger.debug(f"Temporary file deleted: {temp_file_path}")
                except OSError as e:
                    logger.warning(
                        f"Failed to delete temporary file {temp_file_path}: {str(e)}"
                    )
            session.close()

    def _extract_filename_from_response(self, response, url: str) -> str:
        """
        Extract filename from HTTP response (new in v1.10.3).

        Args:
            response: requests HTTP response
            url: Original URL

        Returns:
            str: Filename
        """
        # Try to get filename from Content-Disposition header
        content_disposition = response.headers.get("content-disposition", "")
        if "filename=" in content_disposition:
            import re

            filename_match = re.search(
                r'filename=["\']*([^"\';\r\n]*)', content_disposition
            )
            if filename_match:
                filename = filename_match.group(1).strip()
                if filename:
                    from app.utils import sanitize_filename

                    return sanitize_filename(filename)

        # Use the last segment of the URL as the filename
        from urllib.parse import unquote, urlparse

        parsed_url = urlparse(url)
        filename = unquote(parsed_url.path.split("/")[-1])

        # If there is no extension, try to determine it by Content-Type
        if not get_file_extension(filename):
            content_type = response.headers.get("content-type", "").lower()
            extension = self._get_extension_from_content_type(content_type)
            if extension:
                filename = f"{filename}.{extension}"

        from app.utils import sanitize_filename

        return sanitize_filename(filename) if filename else "downloaded_file"

    def _get_extension_from_content_type(self, content_type: str) -> Optional[str]:
        """
        Determine file extension by Content-Type (new in v1.10.3).

        Args:
            content_type: MIME type

        Returns:
            Optional[str]: File extension or None
        """
        # Mapping of popular MIME types to extensions
        mime_to_extension = settings.MIME_TO_EXTENSION

        # Remove parameters from Content-Type (e.g., charset)
        clean_content_type = content_type.split(";")[0].strip()

        return mime_to_extension.get(clean_content_type)

    def extract_from_url(
        self,
        url: str,
        user_agent: Optional[str] = None,
        extraction_options: Optional[Any] = None,
    ) -> List[Dict[str, Any]]:
        """Extract text from a web page or file by URL (updated in v1.10.3)."""
        # Check URL safety
        if not self._is_safe_url(url):
            raise ValueError(
                "Access to internal IP addresses is prohibited for security reasons"
            )

        try:
            # Step 1: Determine content type via HEAD request
            content_type, final_url = self._determine_content_type(
                url, user_agent, extraction_options
            )

            # Step 2: Choose processing strategy
            if self._is_html_content(content_type, final_url):
                logger.info(
                    f"URL {final_url} determined as HTML page (Content-Type: {content_type}), using web extractor"
                )
                return self._extract_html_page(
                    final_url, user_agent, extraction_options
                )
            else:
                logger.info(
                    f"URL {final_url} determined as file (Content-Type: {content_type}), downloading and processing"
                )
                return self._download_and_extract_file(
                    final_url, user_agent, extraction_options
                )

        except Exception as e:
            logger.error(f"Error processing URL {url}: {str(e)}")
            raise ValueError(f"Error processing URL: {str(e)}")

    def _extract_html_page(
        self,
        url: str,
        user_agent: Optional[str] = None,
        extraction_options: Optional[Any] = None,
    ) -> List[Dict[str, Any]]:
        """
        Extract text from HTML page (separated from extract_from_url in v1.10.3).

        Args:
            url: HTML page URL
            user_agent: Custom User-Agent
            extraction_options: Extraction settings

        Returns:
            List[Dict[str, Any]]: Result of text extraction from page
        """
        html_content = ""
        final_url = url

        # Define settings considering provided parameters or defaults
        enable_javascript = (
            extraction_options.enable_javascript
            if extraction_options and extraction_options.enable_javascript is not None
            else settings.ENABLE_JAVASCRIPT
        )

        # Choose loading method depending on JavaScript settings
        if enable_javascript and sync_playwright:
            logger.info("Using Playwright to load page with JS")
            try:
                html_content, final_url = self._extract_page_with_playwright(
                    url, user_agent, extraction_options
                )
            except Exception as e:
                logger.warning(f"Playwright error: {str(e)}, switching to requests")
                # Fallback to requests on Playwright error
                html_content, final_url = self._extract_page_with_requests(
                    url, user_agent, extraction_options
                )
        else:
            if enable_javascript and not sync_playwright:
                logger.warning(
                    "JavaScript enabled, but Playwright is not installed, using requests"
                )
            logger.info("Using requests to load page")
            html_content, final_url = self._extract_page_with_requests(
                url, user_agent, extraction_options
            )

        try:
            # Extract text from HTML
            page_text = self._extract_text_from_html(html_content)

            # Find and process images
            image_texts = self._extract_images_from_html(
                html_content, final_url, extraction_options
            )

            # Format results
            results = []

            # Add main page content
            results.append(
                {
                    "filename": "page_content",
                    "path": final_url,
                    "size": len(html_content.encode("utf-8")),
                    "type": "html",
                    "text": page_text,
                }
            )

            # Add text from images
            results.extend(image_texts)

            return results

        except Exception as e:
            raise ValueError(f"Error processing HTML page: {str(e)}")

    def _extract_page_with_requests(
        self,
        url: str,
        user_agent: Optional[str] = None,
        extraction_options: Optional[Any] = None,
    ) -> tuple[str, str]:
        """
        Extract HTML content of page using requests (without JS, updated in v1.10.2).

        Args:
            url: Page URL
            user_agent: Custom User-Agent
            extraction_options: Extraction settings

        Returns:
            tuple[str, str]: (html_content, final_url)
        """
        if not requests:
            raise ValueError("requests library not available for web extraction")

        # Define settings considering provided parameters or defaults
        web_page_timeout = (
            extraction_options.web_page_timeout
            if extraction_options and extraction_options.web_page_timeout is not None
            else settings.WEB_PAGE_TIMEOUT
        )

        follow_redirects = (
            extraction_options.follow_redirects
            if extraction_options and extraction_options.follow_redirects is not None
            else True
        )

        # Set User-Agent
        headers = {
            "User-Agent": user_agent or settings.DEFAULT_USER_AGENT,
            "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8",
            "Accept-Language": "ru,en;q=0.5",
            "Accept-Encoding": "gzip, deflate",
            "Connection": "keep-alive",
        }

        try:
            # Load page with timeout
            response = requests.get(
                url,
                headers=headers,
                timeout=web_page_timeout,
                allow_redirects=follow_redirects,
                stream=False,
            )
            response.raise_for_status()

            # Auto-detect encoding
            response.encoding = response.apparent_encoding or "utf-8"
            html_content = response.text
            final_url = response.url

            logger.info(
                f"HTML received via requests, size: {len(html_content)} characters"
            )
            return html_content, final_url

        except requests.RequestException as e:
            if "timeout" in str(e).lower():
                raise ValueError(f"Page loading timeout: {str(e)}")
            elif "connection" in str(e).lower():
                raise ValueError(f"Connection error: {str(e)}")
            else:
                raise ValueError(f"Failed to load page: {str(e)}")

    def _is_safe_url(self, url: str) -> bool:
        """Check URL safety (SSRF protection)."""
        try:
            parsed_url = urlparse(url)

            # Check URL scheme
            if not self._check_url_scheme(parsed_url.scheme):
                return False

            hostname = parsed_url.hostname
            if not hostname:
                logger.warning(f"No hostname in URL: {url}")
                return False

            # Check blocked hostnames
            if not self._check_hostname_not_blocked(hostname, url):
                return False

            # Get host IP addresses
            ips = self._resolve_hostname_ips(hostname)
            if not ips:
                return False

            # Check safety of all IP addresses
            return self._check_all_ips_safe(ips, url)

        except Exception as e:
            logger.warning(f"Error checking URL safety: {str(e)}")
            # Fail-closed: block access in case of error
            return False

    def _check_url_scheme(self, scheme: str) -> bool:
        """Check URL scheme."""
        if scheme not in ["http", "https"]:
            logger.warning(f"Unsupported URL scheme: {scheme}")
            return False
        return True

    def _check_hostname_not_blocked(self, hostname: str, url: str) -> bool:
        """Check that hostname is not blocked."""
        blocked_hostnames = settings.BLOCKED_HOSTNAMES.split(",")
        hostname_lower = hostname.lower()

        for blocked_hostname in blocked_hostnames:
            blocked_hostname = blocked_hostname.strip().lower()
            if blocked_hostname and hostname_lower == blocked_hostname:
                logger.warning(f"Blocked hostname {hostname} for URL {url}")
                return False
        return True

    def _resolve_hostname_ips(self, hostname: str) -> list:
        """Resolve IP addresses for hostname."""
        import socket

        try:
            addr_info = socket.getaddrinfo(
                hostname, None, socket.AF_UNSPEC, socket.SOCK_STREAM
            )
            return [info[4][0] for info in addr_info]
        except socket.gaierror as e:
            logger.warning(f"DNS resolution failed for {hostname}: {str(e)}")
            return []

    def _check_all_ips_safe(self, ips: list, url: str) -> bool:
        """Check safety of all IP addresses."""
        for ip_str in ips:
            if not self._check_single_ip_safe(ip_str, url):
                return False
        return True

    def _check_single_ip_safe(self, ip_str: str, url: str) -> bool:
        """Check safety of a single IP address."""
        try:
            ip_obj = ipaddress.ip_address(ip_str)

            # Check special addresses
            if self._is_special_ip_unsafe(ip_obj, ip_str, url):
                return False

            # Check blocked ranges
            if self._is_ip_in_blocked_ranges(ip_obj, ip_str, url):
                return False

            # Check metadata service
            if self._is_metadata_service_ip(ip_obj, ip_str, url):
                return False

            # Check Docker bridge
            if self._is_docker_bridge_ip(ip_obj, ip_str, url):
                return False

            return True

        except ValueError as e:
            logger.warning(f"Invalid IP address {ip_str}: {str(e)}")
            return True  # Do not block invalid IP

    def _is_special_ip_unsafe(self, ip_obj, ip_str: str, url: str) -> bool:
        """Check for special unsafe IPs."""
        if ip_obj.is_loopback or ip_obj.is_private or ip_obj.is_link_local:
            logger.warning(
                f"Blocked special IP {ip_str} (loopback/private/link-local) for URL {url}"
            )
            return True
        return False

    def _is_ip_in_blocked_ranges(self, ip_obj, ip_str: str, url: str) -> bool:
        """Check if IP belongs to blocked ranges."""
        blocked_ranges = settings.BLOCKED_IP_RANGES.split(",")

        for range_str in blocked_ranges:
            range_str = range_str.strip()
            if not range_str:
                continue

            try:
                network = ipaddress.ip_network(range_str, strict=False)
                if ip_obj in network:
                    logger.warning(
                        f"Blocked IP {ip_str} in range {range_str} for URL {url}"
                    )
                    return True
            except ValueError:
                continue
        return False

    def _is_metadata_service_ip(self, ip_obj, ip_str: str, url: str) -> bool:
        """Check for metadata service IP."""
        if str(ip_obj) == "169.254.169.254":
            logger.warning(f"Blocked metadata service IP {ip_str} for URL {url}")
            return True
        return False

    def _is_docker_bridge_ip(self, ip_obj, ip_str: str, url: str) -> bool:
        """Check for Docker bridge gateway IP."""
        if ip_obj.version == 4:
            octets = str(ip_obj).split(".")
            if (
                octets[0] == "172"
                and 16 <= int(octets[1]) <= 31
                and octets[2] == "0"
                and octets[3] == "1"
            ):
                logger.warning(f"Blocked Docker bridge gateway {ip_str} for URL {url}")
                return True
        return False

    def _extract_text_from_html(self, html_content: str) -> str:
        """Extract text from HTML content."""
        if not BeautifulSoup:
            raise ValueError("BeautifulSoup not available for HTML parsing")

        try:
            soup = BeautifulSoup(html_content, "lxml")

            # Remove scripts, styles and other non-text elements
            for script in soup(["script", "style", "nav", "header", "footer", "aside"]):
                script.decompose()

            # Extract text
            text = soup.get_text()

            # Clean text
            lines = []
            for line in text.splitlines():
                line = line.strip()
                if line:
                    lines.append(line)

            return "\n".join(lines)

        except Exception as e:
            logger.error(f"Error extracting text from HTML: {str(e)}")
            raise ValueError(f"HTML parsing error: {str(e)}")

    def _extract_images_from_html(
        self, html_content: str, base_url: str, extraction_options: Optional[Any] = None
    ) -> List[Dict[str, Any]]:
        """Extract and process images from the page (updated in v1.10.2)."""
        if not BeautifulSoup or not Image:
            return []

        # Setup extraction parameters
        options = self._setup_image_extraction_options(extraction_options)
        if not options["process_images"]:
            logger.info("Image processing disabled in extraction settings")
            return []

        try:
            # Parse images from HTML
            img_tags = self._parse_images_from_html(
                html_content, options["max_images_per_page"]
            )
            if not img_tags:
                return []

            # Categorize images
            base64_images, url_images = self._categorize_images(
                img_tags, options["enable_base64_images"]
            )
            logger.info(
                f"Images found: {len(url_images)} URL, {len(base64_images)} base64"
            )

            results = []
            # Process images
            results.extend(
                self._process_base64_images(base64_images, extraction_options)
            )
            results.extend(
                self._process_url_images(url_images, base_url, extraction_options)
            )

            return results

        except Exception as e:
            logger.warning(f"Error extracting images from HTML: {str(e)}")
            return []

    def _setup_image_extraction_options(
        self, extraction_options: Optional[Any]
    ) -> dict:
        """Setup image extraction parameters."""
        return {
            "process_images": (
                extraction_options.process_images
                if extraction_options and extraction_options.process_images is not None
                else True
            ),
            "max_images_per_page": (
                extraction_options.max_images_per_page
                if extraction_options
                and extraction_options.max_images_per_page is not None
                else settings.MAX_IMAGES_PER_PAGE
            ),
            "enable_base64_images": (
                extraction_options.enable_base64_images
                if extraction_options
                and extraction_options.enable_base64_images is not None
                else settings.ENABLE_BASE64_IMAGES
            ),
        }

    def _parse_images_from_html(self, html_content: str, max_images: int) -> list:
        """Parse images from HTML content."""
        soup = BeautifulSoup(html_content, "lxml")
        img_tags = soup.find_all("img", src=True)
        return img_tags[:max_images]

    def _categorize_images(self, img_tags: list, enable_base64: bool) -> tuple:
        """Categorize images into base64 and URL."""
        base64_images = []
        url_images = []

        for img_tag in img_tags:
            img_src = img_tag.get("src", "")
            if img_src.startswith("data:image/") and enable_base64:
                base64_images.append(img_tag)
            else:
                url_images.append(img_tag)

        return base64_images, url_images

    def _process_base64_images(
        self, base64_images: list, extraction_options: Optional[Any]
    ) -> list:
        """Process base64 images."""
        results = []
        for img_tag in base64_images:
            try:
                result = self._process_base64_image(img_tag, extraction_options)
                if result:
                    results.append(result)
            except Exception as e:
                logger.warning(f"Error processing base64 image: {str(e)}")
        return results

    def _process_url_images(
        self, url_images: list, base_url: str, extraction_options: Optional[Any]
    ) -> list:
        """Process URL images."""
        if not url_images:
            return []

        results = []
        image_download_timeout = (
            extraction_options.image_download_timeout
            if extraction_options
            and extraction_options.image_download_timeout is not None
            else settings.IMAGE_DOWNLOAD_TIMEOUT
        )

        # Process images in groups of 2
        for i in range(0, len(url_images), 2):
            batch = url_images[i : i + 2]
            batch_results = self._process_images_batch(
                batch, base_url, extraction_options, image_download_timeout
            )
            results.extend(batch_results)

        return results

    def _process_images_batch(
        self,
        batch: list,
        base_url: str,
        extraction_options: Optional[Any],
        timeout: int,
    ) -> list:
        """Process a group of images in parallel."""
        batch_results = []

        with concurrent.futures.ThreadPoolExecutor(max_workers=2) as executor:
            futures = []
            for img_tag in batch:
                future = executor.submit(
                    self._process_single_image,
                    img_tag,
                    base_url,
                    extraction_options,
                )
                futures.append(future)

            for future in concurrent.futures.as_completed(futures):
                try:
                    result = future.result(timeout=timeout + 5)
                    if result:
                        batch_results.append(result)
                except Exception as e:
                    logger.warning(f"Error processing image: {str(e)}")

        return batch_results

    def _process_single_image(
        self, img_tag, base_url: str, extraction_options: Optional[Any] = None
    ) -> Optional[Dict[str, Any]]:
        """Process a single image (updated in v1.10.2)."""
        try:
            img_src = img_tag.get("src", "")
            logger.info(f"Processing image: {img_src}")
            if not img_src:
                logger.warning("Image has no src attribute")
                return None

            # Convert relative URL to absolute
            img_url = urljoin(base_url, img_src)
            logger.info(f"Full image URL: {img_url}")

            # Check image URL safety
            if not self._is_safe_url(img_url):
                logger.warning(f"Blocked image URL: {img_url}")
                return None

            # Determine settings
            image_download_timeout = (
                extraction_options.image_download_timeout
                if extraction_options
                and extraction_options.image_download_timeout is not None
                else settings.IMAGE_DOWNLOAD_TIMEOUT
            )

            min_image_size_for_ocr = (
                extraction_options.min_image_size_for_ocr
                if extraction_options
                and extraction_options.min_image_size_for_ocr is not None
                else settings.MIN_IMAGE_SIZE_FOR_OCR
            )

            # Image download
            headers = {"User-Agent": settings.DEFAULT_USER_AGENT, "Referer": base_url}

            response = requests.get(
                img_url, headers=headers, timeout=image_download_timeout, stream=True
            )
            response.raise_for_status()

            # Check image size
            img_content = response.content
            logger.info(f"Image content size: {len(img_content)} bytes")
            if len(img_content) == 0:
                logger.warning("Image content is empty")
                return None

            # Open image to check dimensions
            with Image.open(io.BytesIO(img_content)) as img:
                width, height = img.size
                logger.info(
                    f"Image dimensions: {width}x{height} = {width * height} pixels (min required: {settings.MIN_IMAGE_SIZE_FOR_OCR})"
                )

                # Check minimum size
                if width * height < min_image_size_for_ocr:
                    logger.warning(
                        f"Image too small for OCR: {width * height} < {min_image_size_for_ocr}"
                    )
                    return None

                # OCR image
                logger.info(f"Starting OCR for image: {img_url}")
                text = self._safe_tesseract_ocr(img)
                logger.info(f"OCR result length: {len(text) if text else 0} characters")

                if not text or not text.strip():
                    logger.warning("No text found in image")
                    return None

                # Extract filename from URL
                from .utils import get_extension_from_mime

                filename = os.path.basename(urlparse(img_url).path) or "image"
                if "." not in filename:
                    # Determine extension by MIME type via utility
                    content_type = response.headers.get("content-type", "").lower()
                    extension = get_extension_from_mime(
                        content_type, settings.SUPPORTED_FORMATS
                    )

                    if extension:
                        filename += f".{extension}"
                    else:
                        # If MIME type is not supported, ignore image
                        logger.warning(f"Unsupported image MIME type: {content_type}")
                        return None

                return {
                    "filename": filename,
                    "path": img_url,
                    "size": len(img_content),
                    "type": filename.split(".")[-1].lower(),
                    "text": text.strip(),
                }

        except Exception as e:
            logger.warning(f"Error processing image {img_tag.get('src', '')}: {str(e)}")
            return None

    def _process_base64_image(
        self, img_tag, extraction_options: Optional[Any] = None
    ) -> Optional[Dict[str, Any]]:
        """Process base64 image from data URI (updated in v1.10.2)."""
        try:
            from .utils import (
                decode_base64_image,
                extract_mime_from_base64_data_uri,
                get_extension_from_mime,
            )

            img_src = img_tag.get("src", "")
            logger.info(f"Processing base64 image: {img_src[:50]}...")

            if not img_src.startswith("data:image/"):
                logger.warning("Invalid base64 image format")
                return None

            # Determine settings
            min_image_size_for_ocr = (
                extraction_options.min_image_size_for_ocr
                if extraction_options
                and extraction_options.min_image_size_for_ocr is not None
                else settings.MIN_IMAGE_SIZE_FOR_OCR
            )

            # Extract MIME type
            mime_type = extract_mime_from_base64_data_uri(img_src)
            if not mime_type:
                logger.warning("Could not extract MIME type from base64 image")
                return None

            # Determine file extension
            extension = get_extension_from_mime(mime_type, settings.SUPPORTED_FORMATS)
            if not extension:
                logger.warning(f"Unsupported image MIME type: {mime_type}")
                return None

            # Decode base64 image
            img_content = decode_base64_image(img_src)
            if not img_content:
                logger.warning("Failed to decode base64 image")
                return None

            logger.info(f"Base64 image decoded, size: {len(img_content)} bytes")

            # Open image to check dimensions
            with Image.open(io.BytesIO(img_content)) as img:
                width, height = img.size
                logger.info(
                    f"Base64 image dimensions: {width}x{height} = {width * height} pixels (min required: {min_image_size_for_ocr})"
                )

                # Check minimum size
                if width * height < min_image_size_for_ocr:
                    logger.warning(
                        f"Base64 image too small for OCR: {width * height} < {min_image_size_for_ocr}"
                    )
                    return None

                # OCR изображения
                logger.info("Starting OCR for base64 image")
                text = self._safe_tesseract_ocr(img)
                logger.info(f"OCR result length: {len(text) if text else 0} characters")

                if not text or not text.strip():
                    logger.warning("No text found in base64 image")
                    return None

                # Format filename
                filename = f"base64_image.{extension}"

                return {
                    "filename": filename,
                    "path": f"data:image/{extension};base64,[base64_data]",
                    "size": len(img_content),
                    "type": extension,
                    "text": text.strip(),
                }

        except Exception as e:
            logger.warning(f"Error processing base64 image: {str(e)}")
            return None
